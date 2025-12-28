import os
from typing import Optional
from pypdf import PdfReader
import docx
import openpyxl

def extract_text(path: str) -> Optional[str]:
    ext = os.path.splitext(path)[1].lower()
    try:
        # 텍스트 파일
        if ext in (".txt", ".md", ".markdown", ".log", ".csv", ".json", ".xml", ".html", ".htm", ".css", ".js", ".py", ".java", ".cpp", ".c", ".h"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        
        # PDF 파일
        if ext == ".pdf":
            r = PdfReader(path)
            parts = []
            for p in r.pages:
                t = p.extract_text() or ""
                if t.strip():
                    parts.append(t)
            return "\n\n".join(parts)
        
        # Word 문서
        if ext in (".docx", ".doc"):
            try:
                d = docx.Document(path)
                return "\n".join([p.text for p in d.paragraphs if p.text.strip()])
            except Exception:
                # .doc 파일은 python-docx로 읽을 수 없을 수 있음
                return None
        
        # Excel 파일
        if ext in (".xlsx", ".xls"):
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            parts=[]
            for ws in wb.worksheets:
                parts.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    line = "\t".join("" if v is None else str(v) for v in row)
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts)
        
        # Outlook 파일 (.msg)
        if ext == ".msg":
            try:
                # extract_msg 라이브러리 사용
                try:
                    import extract_msg
                    msg = extract_msg.Message(path)
                    parts = []
                    if msg.subject:
                        parts.append(f"Subject: {msg.subject}")
                    if msg.sender:
                        parts.append(f"From: {msg.sender}")
                    if msg.to:
                        parts.append(f"To: {msg.to}")
                    if msg.cc:
                        parts.append(f"CC: {msg.cc}")
                    if msg.date:
                        parts.append(f"Date: {msg.date}")
                    if msg.body:
                        parts.append(f"\nBody:\n{msg.body}")
                    # 첨부 파일 정보
                    if hasattr(msg, 'attachments') and msg.attachments:
                        parts.append(f"\nAttachments: {len(msg.attachments)}")
                        for att in msg.attachments:
                            if hasattr(att, 'longFilename'):
                                parts.append(f"  - {att.longFilename}")
                    msg.close()
                    return "\n".join(parts)
                except ImportError:
                    # extract_msg가 설치되지 않은 경우 파일명만 반환
                    return f"Outlook message file: {os.path.basename(path)}\nNote: Install 'extract-msg' package to extract content."
            except Exception as e:
                # 오류 발생 시 파일명만 반환
                return f"Outlook message file: {os.path.basename(path)}"
        
        # Outlook 데이터 파일 (.pst, .ost) - 직접 읽기 어려움, 파일명만
        if ext in (".pst", ".ost"):
            return f"Outlook data file: {os.path.basename(path)}\nNote: PST/OST files require special tools to extract content."
        
        # PowerPoint 파일
        if ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(path)
                parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    parts.append(f"# Slide {slide_num}")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                return "\n".join(parts)
            except ImportError:
                return None
            except Exception:
                return None
        
        # 기타 파일 형식은 None 반환 (파일명만 인덱싱됨)
        return None
        
    except Exception as e:
        # 오류 발생 시 None 반환 (파일은 인덱싱되지만 텍스트는 추출 안 됨)
        return None
